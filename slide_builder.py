"""
Slide builder - generate PowerPoint presentations using python-pptx.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from datetime import datetime
from pathlib import Path
import json


# Colors
BG_COLOR = RGBColor(0x1E, 0x1E, 0x2E)  # Dark background
ACCENT_COLOR = RGBColor(0x7C, 0x3A, 0xED)  # Purple
TEXT_COLOR = RGBColor(0xFF, 0xFF, 0xFF)  # White
SUCCESS_COLOR = RGBColor(0x10, 0xB9, 0x81)  # Green
WARNING_COLOR = RGBColor(0xF5, 0x9E, 0x0B)  # Yellow
DANGER_COLOR = RGBColor(0xEF, 0x44, 0x44)  # Red


def add_background(slide, color=BG_COLOR):
    """Add solid background color to slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_slide(prs: Presentation, project: dict, date_str: str):
    """Add title slide with project name, owner, date, and team members."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    add_background(slide)
    
    # Title (project name)
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = project["project_name"]
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle (owner + date)
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(9), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    p = subtitle_frame.paragraphs[0]
    p.text = f"{project['owner']} • {date_str}"
    p.font.size = Pt(20)
    p.font.color.rgb = TEXT_COLOR
    p.alignment = PP_ALIGN.CENTER
    
    # Team members
    team_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(2))
    team_frame = team_box.text_frame
    team_frame.word_wrap = True
    p = team_frame.paragraphs[0]
    p.text = "Team: " + ", ".join(project.get("team_members", []))
    p.font.size = Pt(16)
    p.font.color.rgb = TEXT_COLOR
    p.alignment = PP_ALIGN.CENTER


def add_results_slide(prs: Presentation, results_so_far: str):
    """Add results slide with bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    add_background(slide)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "Results So Far"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8.5), Inches(5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    # Parse bullet points (split by \n or -)
    lines = results_so_far.split("\n")
    for i, line in enumerate(lines):
        if line.strip():
            if i == 0:
                p = content_frame.paragraphs[0]
            else:
                p = content_frame.add_paragraph()
            p.text = line.strip().lstrip("- •")
            p.font.size = Pt(18)
            p.font.color.rgb = TEXT_COLOR
            p.level = 0
            p.space_before = Pt(6)
            p.space_after = Pt(6)


def add_task_status_slide(prs: Presentation, tasks: list):
    """Add task status slide with color-coded table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    add_background(slide)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "Task Status"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR
    
    # Create table
    rows = len(tasks) + 1
    cols = 4
    left = Inches(0.5)
    top = Inches(1.8)
    width = Inches(9)
    height = Inches(4.5)
    
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Header row
    headers = ["Task", "Responsible", "Status", "Due Date"]
    for col_idx, header in enumerate(headers):
        cell = table_shape.cell(0, col_idx)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT_COLOR
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.bold = True
        paragraph.font.color.rgb = BG_COLOR
        paragraph.font.size = Pt(14)
    
    # Data rows
    for row_idx, task in enumerate(tasks):
        status = task.get("status", "not_started")
        
        # Determine row color based on status
        if status == "done":
            row_color = SUCCESS_COLOR
        elif status == "in_progress":
            row_color = WARNING_COLOR
        else:
            row_color = DANGER_COLOR
        
        # Populate cells
        cells_data = [
            task.get("description", "")[:30],
            task.get("responsible", ""),
            status.replace("_", " ").title(),
            task.get("due_date", "")
        ]
        
        for col_idx, cell_text in enumerate(cells_data):
            cell = table_shape.cell(row_idx + 1, col_idx)
            cell.text = cell_text
            cell.fill.solid()
            cell.fill.fore_color.rgb = row_color
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(12)
            paragraph.font.color.rgb = BG_COLOR


def add_next_steps_slide(prs: Presentation, next_steps: str):
    """Add next steps slide with bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    add_background(slide)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "Next Steps"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8.5), Inches(5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    # Parse bullet points
    lines = next_steps.split("\n")
    for i, line in enumerate(lines):
        if line.strip():
            if i == 0:
                p = content_frame.paragraphs[0]
            else:
                p = content_frame.add_paragraph()
            p.text = line.strip().lstrip("- •")
            p.font.size = Pt(18)
            p.font.color.rgb = TEXT_COLOR
            p.level = 0
            p.space_before = Pt(6)
            p.space_after = Pt(6)


def add_timeline_slide(prs: Presentation, timeline: dict, today: str):
    """Add timeline slide with milestone bar."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    add_background(slide)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "Timeline"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR
    
    # Timeline bar background
    start_date = timeline.get("start_date", "")
    end_date = timeline.get("end_date", "")
    
    # Draw timeline background
    bar_left = Inches(1)
    bar_top = Inches(2.5)
    bar_width = Inches(8)
    bar_height = Inches(0.3)
    
    bar = slide.shapes.add_shape(1, bar_left, bar_top, bar_width, bar_height)  # Rectangle
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0x4B, 0x5C, 0x6D)
    bar.line.color.rgb = ACCENT_COLOR
    
    # Add milestones
    milestone_top = Inches(3.2)
    milestones = timeline.get("milestones", [])
    
    for idx, milestone in enumerate(milestones[:5]):  # Limit to 5 milestones
        # Milestone marker
        marker = slide.shapes.add_shape(1, Inches(1.5 + idx * 1.5), bar_top - Inches(0.15), Inches(0.3), Inches(0.6))
        marker.fill.solid()
        marker.fill.fore_color.rgb = ACCENT_COLOR
        marker.line.color.rgb = TEXT_COLOR
        
        # Milestone label
        label_box = slide.shapes.add_textbox(Inches(1.2 + idx * 1.5), milestone_top, Inches(1.2), Inches(1.5))
        label_frame = label_box.text_frame
        label_frame.word_wrap = True
        p = label_frame.paragraphs[0]
        p.text = f"{milestone['label']}\n{milestone['date']}"
        p.font.size = Pt(10)
        p.font.color.rgb = TEXT_COLOR
        p.alignment = PP_ALIGN.CENTER
    
    # Add date range
    date_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(0.6))
    date_frame = date_box.text_frame
    p = date_frame.paragraphs[0]
    p.text = f"{start_date} → {end_date}"
    p.font.size = Pt(16)
    p.font.color.rgb = TEXT_COLOR
    p.alignment = PP_ALIGN.CENTER


def generate_slides(project: dict, results_so_far: str, next_steps: str) -> str:
    """
    Generate a 5-slide PowerPoint presentation.
    
    Args:
        project: Project dict
        results_so_far: Results text
        next_steps: Next steps text
        
    Returns:
        Path to the generated .pptx file
    """
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    today = datetime.now().strftime("%Y-%m-%d")
    
    # Generate slides
    add_title_slide(prs, project, today)
    add_results_slide(prs, results_so_far)
    add_task_status_slide(prs, project.get("tasks", []))
    add_next_steps_slide(prs, next_steps)
    add_timeline_slide(prs, project.get("timeline", {}), today)
    
    # Save file
    slides_dir = Path(f"slides/{project['project_name']}")
    slides_dir.mkdir(parents=True, exist_ok=True)
    
    filename = f"{project['project_name']}_{today}.pptx"
    filepath = slides_dir / filename
    
    prs.save(str(filepath))
    return str(filepath)
